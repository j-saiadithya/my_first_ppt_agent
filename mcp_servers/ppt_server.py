"""
PPT MCP Server — Step 2 & 3
============================
This MCP server exposes PowerPoint creation tools via the Model Context Protocol.

Tools provided:
  • create_presentation(title)   → Start a new presentation with a title slide
  • add_slide(title, bullets)    → Add a content slide with bullet points
  • save_presentation(filename)  → Save the .pptx file to the output/ folder
  • get_slide_count()            → Return how many slides exist (for observation)

The agent communicates with this server over stdio using the MCP protocol.
It should NEVER import python-pptx directly — all slide work happens here.
"""

import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from mcp.server.fastmcp import FastMCP

# ---------------------------------------------------------------------------
# MCP Server Instance
# ---------------------------------------------------------------------------
mcp = FastMCP("ppt-server")

# ---------------------------------------------------------------------------
# Internal State — holds the active presentation in memory
# ---------------------------------------------------------------------------
_state = {
    "presentation": None,
    "title": None,
}

# ---------------------------------------------------------------------------
# Helper: resolve output directory (relative to project root)
# ---------------------------------------------------------------------------
PROJECT_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
OUTPUT_DIR = os.path.join(PROJECT_ROOT, "output")


def _ensure_output_dir():
    """Create the output/ folder if it doesn't exist."""
    os.makedirs(OUTPUT_DIR, exist_ok=True)


# ---------------------------------------------------------------------------
# Tool 1: create_presentation
# ---------------------------------------------------------------------------
@mcp.tool()
def create_presentation(title: str) -> str:
    """
    Create a new PowerPoint presentation with a title slide.

    Args:
        title: The presentation title displayed on the first slide.

    Returns:
        A confirmation message.
    """
    prs = Presentation()

    # --- Title slide ---
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title text
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Style the title
    for paragraph in title_placeholder.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)  # Dark navy

    # Set subtitle if available
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "Auto-generated presentation"
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Store in state
    _state["presentation"] = prs
    _state["title"] = title

    return json.dumps({
        "status": "success",
        "message": f"Presentation '{title}' created with title slide.",
        "slide_count": len(prs.slides),
    })


# ---------------------------------------------------------------------------
# Tool 2: add_slide
# ---------------------------------------------------------------------------
@mcp.tool()
def add_slide(title: str, bullets: str) -> str:
    """
    Add a content slide with a title and bullet points.

    Args:
        title: The slide title.
        bullets: A JSON-encoded list of bullet point strings.
                 Example: '["Point 1", "Point 2", "Point 3"]'

    Returns:
        A confirmation message with the current slide count.
    """
    if _state["presentation"] is None:
        return json.dumps({
            "status": "error",
            "message": "No presentation exists. Call create_presentation first.",
        })

    prs = _state["presentation"]

    # Parse bullets — accept JSON list or comma-separated fallback
    try:
        bullet_list = json.loads(bullets)
        if not isinstance(bullet_list, list):
            bullet_list = [str(bullet_list)]
    except (json.JSONDecodeError, TypeError):
        # Fallback: split by newlines or commas
        bullet_list = [b.strip() for b in str(bullets).replace("\\n", "\n").split("\n") if b.strip()]
        if len(bullet_list) <= 1:
            bullet_list = [b.strip() for b in str(bullets).split(",") if b.strip()]

    # Validate: 3-5 bullets, ≤12 words each
    # Truncate if too many
    if len(bullet_list) > 5:
        bullet_list = bullet_list[:5]

    # Pad if too few
    while len(bullet_list) < 3:
        bullet_list.append("Additional details to be added")

    # Truncate long bullets to ~12 words
    cleaned_bullets = []
    for bullet in bullet_list:
        words = str(bullet).split()
        if len(words) > 12:
            words = words[:12]
        cleaned_bullets.append(" ".join(words))

    # --- Create the slide ---
    slide_layout = prs.slide_layouts[1]  # Title + Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    slide_title = slide.shapes.title
    slide_title.text = title
    for paragraph in slide_title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Set bullet points in the content placeholder
    content_placeholder = slide.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()

    for i, bullet in enumerate(cleaned_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = bullet
        p.level = 0
        p.space_after = Pt(8)

        for run in p.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    return json.dumps({
        "status": "success",
        "message": f"Slide '{title}' added with {len(cleaned_bullets)} bullet points.",
        "slide_count": len(prs.slides),
    })


# ---------------------------------------------------------------------------
# Tool 3: save_presentation
# ---------------------------------------------------------------------------
@mcp.tool()
def save_presentation(filename: str) -> str:
    """
    Save the current presentation to a .pptx file in the output/ folder.

    Args:
        filename: The filename (without path). Example: 'my_presentation.pptx'

    Returns:
        The full path to the saved file.
    """
    if _state["presentation"] is None:
        return json.dumps({
            "status": "error",
            "message": "No presentation to save. Call create_presentation first.",
        })

    # Ensure .pptx extension
    if not filename.endswith(".pptx"):
        filename += ".pptx"

    # Sanitize filename (remove characters that are invalid on Windows)
    safe_name = "".join(c for c in filename if c not in r'<>:"/\|?*').strip()
    if not safe_name:
        safe_name = "presentation.pptx"

    _ensure_output_dir()
    filepath = os.path.join(OUTPUT_DIR, safe_name)

    _state["presentation"].save(filepath)

    # Reset state after saving
    saved_title = _state["title"]
    _state["presentation"] = None
    _state["title"] = None

    return json.dumps({
        "status": "success",
        "message": f"Presentation '{saved_title}' saved successfully.",
        "filepath": filepath,
    })


# ---------------------------------------------------------------------------
# Tool 4: get_slide_count
# ---------------------------------------------------------------------------
@mcp.tool()
def get_slide_count() -> str:
    """
    Get the current number of slides in the presentation.

    Returns:
        The slide count or an error if no presentation exists.
    """
    if _state["presentation"] is None:
        return json.dumps({
            "status": "error",
            "message": "No active presentation.",
            "slide_count": 0,
        })

    count = len(_state["presentation"].slides)
    return json.dumps({
        "status": "success",
        "slide_count": count,
        "title": _state["title"],
    })


# ---------------------------------------------------------------------------
# Run the MCP server (stdio transport)
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    mcp.run()
